import pandas as pd
import re
import os
import json
from PIL import Image
from logging import PlaceHolder
from pptx import Presentation
from pptx.util import Pt

with open("config.json") as json_file:
    data = json.load(json_file)

    NameCol = data["column"]['actualName']
    DisplayCol = data["column"]['displayName']
    YearCol = data["column"]['year']
    MajorCol = data["column"]['major']
    CaptionCol = data["column"]['caption']

    ExcelLocation = data["location"]["excel"]
    TemplateLocation = data["location"]["template"]
    FontLocation = data["location"]["font"]
    PhotoLocation = data["location"]["photo"]
    PptxDestination = data["location"]["target"]

def ProcessField(s):
    return re.sub(r'(?<!\w)and(?!\w)',"&", s).upper()

def GetFileName(name):
    name = name[:20].strip()
    picArr = os.listdir(PhotoLocation)
    filename = []
    for x in picArr:
        searched = re.search(name, x)
        if searched:
            filename.append(x)
    if len(filename) == 0:
        print("No picture found for {0}".format(name))
        return None
    if len(filename) > 1:
        print("Ambiguity in filename for {0}. Choosing first path".format(name), filename)
    return filename[0]

def PrimePics():
    imagenames = os.listdir(PhotoLocation)
    for imagename in imagenames:
        path = os.path.join(PhotoLocation, imagename)
        try:
            im = Image.open(path)
            im.save(path)
        except Exception as e:
            print(e)

def CreateDoorcard(name, data_dict):
    prs = Presentation(TemplateLocation)
    phs = prs.slides[0].placeholders
    for ph in phs:
        if ph.name == "Picture":
            ph.insert_picture(os.path.join(PhotoLocation, GetFileName(name)))
        else:
            ph.text = ProcessField(data_dict[ph.name])

    safeName = re.sub(r'[^A-z]', "", name) + "_Ursa.pptx"
    if not os.path.exists(PptxDestination):
        os.makedirs(PptxDestination)
    prs.save(os.path.join(PptxDestination, safeName))

if __name__ == "__main__":
    df = pd.read_excel(ExcelLocation)
    PrimePics()
    for i, row in df.iterrows():
        try:
            CreateDoorcard(
                row[NameCol], 
                {
                    "Name": row[DisplayCol], 
                    "Year": row[YearCol], 
                    "Major": row[MajorCol], 
                    "Caption": row[CaptionCol]
                }
            )
        except Exception as e:
            print("Error at {0}:".format(row[NameCol]), e)